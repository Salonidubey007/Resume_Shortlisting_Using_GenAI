# utils/file_reader.py

from pathlib import Path
import subprocess
from utils.logger import log_warn

try:
    from docx import Document
except:
    Document = None

try:
    from PyPDF2 import PdfReader
except:
    PdfReader = None

try:
    from pptx import Presentation
except:
    Presentation = None


SUPPORTED_EXTS = {".pdf", ".docx", ".doc", ".txt", ".pptx"}


def read_txt(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except:
        return path.read_text(encoding="latin-1", errors="ignore")


def read_docx(path: Path) -> str:
    if Document is None:
        raise RuntimeError("Missing python-docx")

    doc = Document(str(path))
    parts = []

    for p in doc.paragraphs:
        text = (p.text or "").strip()
        if text:
            parts.append(text)

    for ti, table in enumerate(doc.tables):
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def read_pdf(path: Path) -> str:
    if PdfReader is None:
        raise RuntimeError("Missing PyPDF2")

    try:
        reader = PdfReader(str(path))
    except Exception as e:
        log_warn(f"[PDF] Failed reading {path.name}: {e}")
        return ""

    pages = []
    for page in reader.pages:
        try:
            t = page.extract_text() or ""
        except:
            t = ""
        if t.strip():
            pages.append(t)

    return "\n".join(pages)


def read_pptx(path: Path) -> str:
    if Presentation is None:
        raise RuntimeError("Missing python-pptx")

    prs = Presentation(str(path))
    parts = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    parts.append(shape.text.strip())

    return "\n".join(parts)


def convert_doc_to_docx(doc_path: Path, out_dir: Path) -> Path:
    out_dir.mkdir(parents=True, exist_ok=True)
    cmd = ["soffice", "--headless", "--convert-to", "docx", "--outdir", str(out_dir), str(doc_path)]

    try:
        subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
    except Exception as e:
        raise RuntimeError(f"[DOC] Conversion failed: {e}")

    out_path = out_dir / f"{doc_path.stem}.docx"
    return out_path


def read_any(path: Path, tmp_dir: Path = Path(".tmp_doc_convert")) -> str:
    ext = path.suffix.lower()

    if ext == ".txt":
        return read_txt(path)

    if ext == ".pdf":
        return read_pdf(path)

    if ext == ".docx":
        return read_docx(path)

    if ext == ".pptx":
        return read_pptx(path)

    if ext == ".doc":
        docx_path = convert_doc_to_docx(path, tmp_dir)
        return read_docx(docx_path)

    raise ValueError(f"Unsupported file: {path}")